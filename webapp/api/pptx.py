from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

# Create a new PowerPoint presentation
prs = Presentation()
slide_layout = prs.slide_layouts[5]  # Title Only Layout
slide = prs.slides.add_slide(slide_layout)

# Add title to the slide
title = slide.shapes.title
title.text = "Relevant Literature Review"

# Define table dimensions
rows, cols = 3, 4
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(2)

# Add table to slide
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
col_widths = [Inches(2), Inches(3.5), Inches(1.5), Inches(2)]
for i, width in enumerate(col_widths):
    table.columns[i].width = width

# Add header row with formatting
headers = ["Author’s Name", "Paper name and publication details", "Year of publication", "Main content of the paper"]
for col, header_text in enumerate(headers):
    cell = table.cell(0, col)
    cell.text = header_text
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(207, 83, 0)  # Orange background
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.bold = True
        paragraph.font.size = Inches(0.3)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Add data to the table
data = [
    ["Pérez-Molina", "Exploring a multilevel approach with spatial effects to model housing price", "2021",
     "A multilevel approach to model housing price"],
    ["Griffith", "Spatially varying coefficient models in real estate", "2016",
     "Eigenvector spatial filtering and alternative approaches"]
]

for row_idx, row_data in enumerate(data, start=1):
    for col_idx, text in enumerate(row_data):
        table.cell(row_idx, col_idx).text = text

# Save the presentation file
pptx_path = "/mnt/data/Literature_Review_Slide.pptx"
prs.save(pptx_path)
pptx_path
